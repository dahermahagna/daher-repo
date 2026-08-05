#!/usr/bin/env python3
"""Generate BIOS Auto-Update design presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x35, 0x64)   # slide background / headings
MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # accent / shapes
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)   # subtle fill
ORANGE      = RGBColor(0xED, 0x7D, 0x31)   # highlight
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
GREEN       = RGBColor(0x70, 0xAD, 0x47)
RED         = RGBColor(0xFF, 0x4B, 0x4B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid() if fill_color else shape.fill.background()
    if fill_color:
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def label(slide, text, l, t, w, h,
          font_size=Pt(14), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def heading(slide, text, sub=None):
    """Full-width light header bar."""
    box(slide, 0, 0, 13.33, 1.1, fill_color=LIGHT_BLUE)
    label(slide, text, 0.3, 0.15, 12.5, 0.6,
          font_size=Pt(28), bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)
    if sub:
        label(slide, sub, 0.3, 0.7, 12.5, 0.35,
              font_size=Pt(14), color=MID_BLUE, align=PP_ALIGN.LEFT)


def arrow(slide, x1, y1, x2, y2, color=MID_BLUE):
    """Draw a simple horizontal or vertical connector arrow."""
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    from lxml import etree

    sp = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    sp.line.color.rgb = color
    sp.line.width = Pt(1.5)
    # add arrowhead via XML
    ln = sp._element.find(qn('p:spPr') + '/' + '{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if ln is None:
        ln = sp.line._ln
    tail = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail.set('type', 'arrow')


def card(slide, l, t, w, h, title, body_lines,
         title_color=DARK_BLUE, body_color=DARK_GREY,
         fill=WHITE, border=MID_BLUE):
    box(slide, l, t, w, h, fill_color=fill, line_color=border, line_width=Pt(1.5))
    label(slide, title, l + 0.1, t + 0.05, w - 0.2, 0.35,
          font_size=Pt(12), bold=True, color=title_color)
    body = "\n".join(body_lines)
    label(slide, body, l + 0.1, t + 0.38, w - 0.2, h - 0.45,
          font_size=Pt(10), color=body_color)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)

# accent stripe
box(s, 0, 3.2, 13.33, 0.08, fill_color=ORANGE)

label(s, "BIOS Auto-Update System", 0.8, 1.4, 11.7, 1.2,
    font_size=Pt(44), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
label(s, "Design Overview", 0.8, 2.6, 11.7, 0.7,
    font_size=Pt(26), bold=False, color=MID_BLUE, align=PP_ALIGN.CENTER)
label(s, "Multi-jump firmware upgrade with resume-on-reboot",
      0.8, 3.5, 11.7, 0.5,
      font_size=Pt(18), color=RGBColor(0xBF, 0xBF, 0xBF), align=PP_ALIGN.CENTER)
label(s, "bios_update_tools.py  |  bios-update.py  |  config.py",
      0.8, 6.6, 11.7, 0.5,
      font_size=Pt(12), color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 – Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Problem Statement",
        "Why a plain firmware flash is not enough")

problems = [
    ("1  Multi-jump upgrades",
     "BIOS cannot always jump directly to the target version.\n"
     "Intermediate versions must be applied in strict order."),
    ("2  Reboot required between jumps",
     "Each jump triggers a reboot. Progress must survive power cycles\n"
     "so the next boot resumes from the correct step."),
    ("3  Retry on transient failure",
     "Flash operations can fail. Each step must be retried up to\n"
     "N times before the whole plan is marked failed."),
    ("4  Base plan must stay immutable",
     "The operator-provided plan file must never be modified.\n"
     "Only a separate runtime state file tracks mutable progress."),
]

col_w, col_h = 5.9, 2.2
positions = [(0.3, 1.4), (7.1, 1.4), (0.3, 3.8), (7.1, 3.8)]
for (title, body), (l, t) in zip(problems, positions):
    card(s, l, t, col_w, col_h, title, body.split("\n"),
         fill=WHITE, border=MID_BLUE, title_color=MID_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 – Component Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Component Architecture", "Files and their responsibilities")

# config.py
card(s, 0.3, 1.3, 3.5, 2.1, "config.py",
     ["VERSION_JUMPS  – jump table (from→to→file)",
      "BIOSVersionsDict  – current target per board",
      "PAYLOAD_BASE_PATH  – firmware directory"],
     fill=LIGHT_BLUE, border=DARK_BLUE, title_color=DARK_BLUE)

# bios_update_tools.py
card(s, 4.9, 1.3, 3.5, 2.1, "bios_update_tools.py",
    ["load_upgrade_plan_with_state + save_upgrade_state",
      "PlanValidationError  – typed error with code",
      "Helper functions (exec, board, platform…)"],
     fill=LIGHT_BLUE, border=DARK_BLUE, title_color=DARK_BLUE)

# bios-update.py
card(s, 9.5, 1.3, 3.5, 2.1, "bios-update.py",
     ["CLI entry point",
    "parse_args → load / simulate",
    "save runtime state to disk"],
     fill=LIGHT_BLUE, border=DARK_BLUE, title_color=DARK_BLUE)

# arrows config → tools → cli
arrow(s, 3.8, 2.35, 4.9, 2.35)
arrow(s, 8.4, 2.35, 9.5, 2.35)

# File boxes
card(s, 0.3, 4.0, 3.9, 1.2, "base_plan.json  (read-only)",
     ["Operator-provided step definitions",
      "from_version / to_version / payload_path"],
     fill=WHITE, border=ORANGE, title_color=ORANGE)

card(s, 4.7, 4.0, 3.9, 1.2, "upgrade_plan_state.json  (mutable)",
     ["Runtime state: status, attempts, current ids",
      "Written after each step; base plan untouched"],
     fill=WHITE, border=GREEN, title_color=GREEN)

card(s, 9.1, 4.0, 3.9, 1.2, "runtime_state.json  (optional override path)",
    ["Defaults to /var/lib/harmonic-bios-autoupdate path",
     "Can be provided via --state-file"],
    fill=WHITE, border=MID_BLUE, title_color=MID_BLUE)

# tools ↔ files arrows
arrow(s, 4.9, 3.4, 2.25, 4.0)
arrow(s, 6.65, 3.4, 6.65, 4.0)
arrow(s, 8.4, 3.4, 11.05, 4.0)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 – Execution Path
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Execution Path",
      "The CLI runs in user-provided base-plan mode")

box(s, 3.6, 1.25, 6.1, 5.7, fill_color=WHITE, line_color=ORANGE, line_width=Pt(2))
label(s, "Base Plan + State (load/execute/resume)",
    3.8, 1.3, 5.7, 0.4, font_size=Pt(13), bold=True, color=ORANGE)

steps_b = [
    "① parse --base-plan-file + --state-file",
    "② load read-only base plan JSON",
    "③ load state file  OR  build default state",
    "④ validate_upgrade_state()",
    "⑤ merge: overlay state onto base plan",
    "⑥ save_upgrade_state() → state file only",
]
for i, step in enumerate(steps_b):
    box(s, 3.8, 1.85 + i * 0.72, 5.7, 0.6,
        fill_color=RGBColor(0xFF, 0xF0, 0xE0), line_color=None)
    label(s, step, 3.95, 1.88 + i * 0.72, 5.4, 0.55,
          font_size=Pt(11), color=DARK_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 – Step Execution & Retry Logic
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Step Execution & Retry Logic",
        "simulate_plan_transitions() – how each step is driven")

# flowchart drawn with boxes + arrows
def fbox(l, t, w, h, text, fc, tc=DARK_GREY, fs=Pt(11)):
    box(s, l, t, w, h, fill_color=fc, line_color=DARK_BLUE, line_width=Pt(1))
    label(s, text, l + 0.08, t + 0.05, w - 0.16, h - 0.1,
          font_size=fs, color=tc, align=PP_ALIGN.CENTER, wrap=True)

cx = 5.0   # center x
fbox(cx - 1.4, 1.2, 2.8, 0.55, "state = in_progress", LIGHT_BLUE, fs=Pt(12))
arrow(s, cx, 1.75, cx, 2.05)
fbox(cx - 1.4, 2.05, 2.8, 0.55, "next step: status = in_progress", LIGHT_BLUE, fs=Pt(11))
arrow(s, cx, 2.6, cx, 2.9)
fbox(cx - 1.4, 2.9, 2.8, 0.55, "attempt += 1  →  apply outcome", WHITE, fs=Pt(11))
arrow(s, cx, 3.45, cx, 3.75)

# decision diamond (approximated as rotated rectangle via shape)
fbox(cx - 1.4, 3.75, 2.8, 0.6, "outcome == success?", LIGHT_BLUE, fs=Pt(12))

# YES branch (right)
arrow(s, cx + 1.4, 4.05, cx + 2.8, 4.05)
label(s, "YES", cx + 1.45, 3.85, 1.3, 0.3, font_size=Pt(10), color=GREEN, bold=True)
fbox(cx + 2.8, 3.75, 2.5, 0.6, "step.status = done", GREEN, tc=WHITE, fs=Pt(11))
arrow(s, cx + 4.05, 4.35, cx + 4.05, 4.65)
fbox(cx + 2.8, 4.65, 2.5, 0.6, "more steps?  → repeat", LIGHT_BLUE, fs=Pt(11))
arrow(s, cx + 4.05, 5.25, cx + 4.05, 5.55)
fbox(cx + 2.8, 5.55, 2.5, 0.6, "state = completed ✓", GREEN, tc=WHITE, fs=Pt(12))

# NO branch (left)
arrow(s, cx - 1.4, 4.05, cx - 2.8, 4.05)
label(s, "NO", cx - 2.7, 3.85, 0.8, 0.3, font_size=Pt(10), color=RED, bold=True)
fbox(cx - 5.3, 3.75, 2.5, 0.6, "attempts < max?", LIGHT_BLUE, fs=Pt(11))
arrow(s, cx - 4.05, 3.75, cx - 4.05, 3.45)
label(s, "YES → retry", cx - 5.0, 3.2, 2.0, 0.3, font_size=Pt(10), color=DARK_GREY)
arrow(s, cx - 5.3, 4.05, cx - 6.5, 4.05)
label(s, "NO", cx - 6.45, 3.85, 0.6, 0.3, font_size=Pt(10), color=RED, bold=True)
fbox(cx - 8.5, 3.75, 2.0, 0.6, "step.status = failed\nstate = failed ✗", RED, tc=WHITE, fs=Pt(10))

# note bottom
label(s, "max_attempts_per_step  (default: 2)  is set in the plan and validated on load",
      0.5, 6.8, 12.3, 0.4,
      font_size=Pt(11), color=DARK_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 – State File: Resume Across Reboots
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Resume Across Reboots",
        "The base plan is never modified — only the state file changes")

# Timeline: 3 boots
boot_y = 2.8
for i, (bx, label_text) in enumerate([(1.5, "Boot 1"), (6.0, "Boot 2"), (10.5, "Boot 3")]):
    box(s, bx - 0.55, boot_y - 0.3, 1.1, 0.6,
        fill_color=DARK_BLUE, line_color=None)
    label(s, label_text, bx - 0.5, boot_y - 0.25, 1.0, 0.5,
          font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# connecting line
box(s, 0.9, boot_y + 0.15, 11.6, 0.04, fill_color=MID_BLUE)

# Step boxes per boot
boot_details = [
    [("load base plan", MID_BLUE), ("build default state", MID_BLUE),
     ("execute step 1", GREEN), ("save state (step1=done)", ORANGE)],
    [("load base plan", MID_BLUE), ("load state (step1=done)", ORANGE),
     ("execute step 2", GREEN), ("save state (step2=done)", ORANGE)],
    [("load base plan", MID_BLUE), ("load state (step2=done)", ORANGE),
     ("execute step 3", GREEN), ("state = completed ✓", GREEN)],
]
col_starts = [0.3, 4.8, 9.3]
for col, (details, cx) in enumerate(zip(boot_details, col_starts)):
    for row, (txt, color) in enumerate(details):
        box(s, cx, 3.7 + row * 0.72, 3.9, 0.62,
            fill_color=color if color != GREEN else RGBColor(0xE2, 0xEF, 0xDA),
            line_color=color, line_width=Pt(1))
        tc = DARK_GREY if color != DARK_BLUE else WHITE
        label(s, txt, cx + 0.08, 3.73 + row * 0.72, 3.74, 0.55,
              font_size=Pt(10), color=tc)

label(s, "★  Base plan file on disk is identical before and after all three boots",
      0.5, 6.85, 12.3, 0.4,
      font_size=Pt(12), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 – Validation Layers
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Validation Layers",
        "Every data boundary is validated before use")

layers = [
    ("validate_base_upgrade_plan()",
     ["Required top-level fields", "step_id / from / to / payload per step",
      "action_id per action"],
     MID_BLUE),
    ("validate_upgrade_state()",
     ["step count matches base plan", "status ∈ {pending,in_progress,done,failed}",
      "attempts ≥ 0, current_step_index in range"],
     ORANGE),
    ("validate_upgrade_plan()",
     ["Full runtime plan after merge", "All step fields present",
      "current_step_index ≤ len(steps)"],
     GREEN),
]

col_w, col_h = 5.9, 2.5
positions = [(0.3, 1.4), (7.1, 1.4), (0.3, 4.1), (7.1, 4.1)]
for (fn, checks, color), (l, t) in zip(layers, positions):
    box(s, l, t, col_w, col_h, fill_color=WHITE, line_color=color, line_width=Pt(2))
    label(s, fn, l + 0.15, t + 0.1, col_w - 0.3, 0.4,
          font_size=Pt(13), bold=True, color=color)
    for i, check in enumerate(checks):
        label(s, "✔  " + check, l + 0.15, t + 0.6 + i * 0.55, col_w - 0.3, 0.5,
              font_size=Pt(11), color=DARK_GREY)

label(s, "PlanValidationError(code, message)  is raised on any failure — callers catch by .code",
      0.5, 6.85, 12.3, 0.4,
      font_size=Pt(12), bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

label(s, "Example schema error code: MISSING_BASE_ACTION_COMMAND",
    0.5, 6.45, 12.3, 0.35,
    font_size=Pt(10), color=MID_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 – Key Design Decisions
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s, LIGHT_GREY)
heading(s, "Key Design Decisions")

decisions = [
    ("Base plan is immutable",
     "Operator intent is never overwritten.  Only upgrade_plan_state.json changes."),
    ("State file separation",
     "Allows safe resume after any reboot or crash without re-reading operator config."),
    ("Typed PlanValidationError",
     "Carries a machine-readable .code so callers and tests can match specific failures."),
    ("raise_on_error toggle",
     "Validation can either raise (strict / test mode) or log-and-return-None (service mode)."),
    ("max_attempts_per_step in state",
     "Retry budget is part of persisted state, not hard-coded, so it survives across boots."),
    ("Single execution path",
     "Only user-provided base-plan path is supported by CLI and runtime manager."),
]

for i, (title, body) in enumerate(decisions):
    row = i % 3
    col = i // 3
    l = 0.4 + col * 6.5
    t = 1.4 + row * 1.85
    box(s, l, t, 6.1, 1.65, fill_color=WHITE, line_color=ORANGE, line_width=Pt(1.5))
    label(s, title, l + 0.15, t + 0.1, 5.8, 0.45,
          font_size=Pt(13), bold=True, color=ORANGE)
    label(s, body, l + 0.15, t + 0.55, 5.8, 1.0,
            font_size=Pt(11), color=DARK_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
out = "/home/dmahagne/test-dir/python_tries/bios_update/bios_update_design.pptx"
prs.save(out)
print(f"Saved: {out}")
